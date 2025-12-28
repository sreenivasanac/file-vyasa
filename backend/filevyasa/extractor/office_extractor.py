"""Office document extractor for Word, Excel, PowerPoint, and OpenDocument formats."""

from pathlib import Path
from typing import Any, Dict, Tuple

import structlog

from filevyasa.extractor.base import BaseExtractor, MarkItDownMixin

logger = structlog.get_logger()

# Minimum text length to consider extraction successful
MIN_TEXT_LENGTH = 10


class OfficeExtractor(MarkItDownMixin, BaseExtractor):
    """Extractor for Office documents using markitdown with python-docx fallback."""

    @classmethod
    def supported_extensions(cls) -> list[str]:
        return [
            # Microsoft Office
            "docx", "doc",
            "xlsx", "xls", "csv", "tsv",
            "pptx", "ppt",
            # Rich Text
            "rtf",
            # OpenDocument
            "odt", "ods", "odp",
        ]

    def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from Office documents."""
        metadata = {}
        ext = file_path.suffix.lower()

        # Try markitdown first
        md = self._get_markitdown()
        if md is not None:
            try:
                result = md.convert(str(file_path))
                content = result.text_content if result.text_content else ""

                if len(content.strip()) >= MIN_TEXT_LENGTH:
                    metadata["extraction_method"] = "markitdown"
                    metadata["source_type"] = ext
                    return content, metadata

                logger.info(
                    "markitdown_insufficient_content",
                    path=str(file_path),
                    text_length=len(content.strip()),
                )
            except Exception as e:
                logger.error("markitdown_extraction_failed", path=str(file_path), error=str(e))

        # Fall back to type-specific extractors
        return self._fallback_extract(file_path, metadata)

    def _fallback_extract(self, file_path: Path, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
        """Fallback extraction for when markitdown fails."""
        metadata["extraction_method"] = "fallback"
        ext = file_path.suffix.lower()

        # CSV/TSV can be read as text
        if ext in (".csv", ".tsv"):
            try:
                content = file_path.read_text(encoding="utf-8")
                metadata["extraction_method"] = "text_read"
                return content, metadata
            except Exception:
                pass

        # For DOCX, try python-docx
        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(str(file_path))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                content = "\n\n".join(paragraphs)
                metadata["paragraph_count"] = len(paragraphs)
                metadata["extraction_method"] = "python-docx"
                return content, metadata
            except Exception as e:
                logger.warning("python_docx_failed", error=str(e))

        # For XLSX, try openpyxl
        if ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(file_path), data_only=True)
                rows = []
                for sheet in wb.worksheets:
                    rows.append(f"## {sheet.title}")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join(str(c) if c is not None else "" for c in row)
                        if row_text.strip(" |"):
                            rows.append(row_text)
                content = "\n".join(rows)
                metadata["sheet_count"] = len(wb.worksheets)
                metadata["extraction_method"] = "openpyxl"
                return content, metadata
            except Exception as e:
                logger.warning("openpyxl_failed", error=str(e))

        # For PPTX, try python-pptx
        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                texts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text)
                    if slide_texts:
                        texts.append(f"## Slide {slide_num}\n" + "\n".join(slide_texts))
                content = "\n\n".join(texts)
                metadata["slide_count"] = len(prs.slides)
                metadata["extraction_method"] = "python-pptx"
                return content, metadata
            except Exception as e:
                logger.warning("python_pptx_failed", error=str(e))

        return f"[Unable to extract content from {ext} file]", metadata
